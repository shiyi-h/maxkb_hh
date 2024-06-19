#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
from io import BytesIO
from pptx import Presentation
from PIL import Image
from deepdoc.vision import OCR
import re
from collections import OrderedDict
import numpy as np
from dataset.models import Image as ImageModel
import uuid
from common.handle.impl.utils import emf_to_png
import hashlib



def html_table(table):
    tb = '<table><tr>'
    for value in table[0]:
        tb += f"<th>{value}</th>"
    tb += "</tr>"
    for row in table[1:]:
        tb += "<tr>"
        for value in row:
            tb += f"<td>{value}</td>"
        tb += "</tr>"
    tb += "</table>\n"
    return tb 

def md_table(table):
    tb = '| ' + ' | '.join([value for value in table[0]]) + ' |\n'
    tb += '| ' + ' | '.join(['---' for i in range(len(table[0]))]) + ' |\n'
    for row in table[1:]:
        tb += '| ' + ' | '.join([value for value in row]) + ' |\n'
    return tb

def build_table(text):
    if not isinstance(text, str):
        text = '\n'.join(text)
    # 匹配##{top},{left}##格式的坐标并提取坐标和内容
    pattern = re.compile(r'(.*?)@@(-?\d+),(-?\d+)##')
    matches = pattern.findall(text)
    
    # 以左坐标left为键，存储每列的所有(top, content)元组
    columns = OrderedDict()
    for content, top, left in matches:
        left = int(left)
        top = int(top)
        if left not in columns:
            columns[left] = []
        columns[left].append((top, content))
    if not columns:
        return text
    
    # 确定每列的最大行数，用于验证是否可以组成表格
    rows_counter = {}
    for col in columns.values():
        rows = tuple(sorted([row for row,_ in col]))
        rows_counter[rows] = rows_counter.get(rows, 0)+1
        
    tables_dict = {}
    for col in columns:
        rows = tuple(sorted([row for row,_ in columns[col]]))
        if len(rows) == 1:
            continue
        if rows_counter[rows]>1:
            for row, content in columns[col]:
                if rows in tables_dict:
                    if row in tables_dict[rows]:
                        tables_dict[rows][row][col] = content
                    else:
                        tables_dict[rows][row] = {col: content}
                else:
                    tables_dict[rows] = {row:{col:content}}
                text = text.replace(f'{content}@@{row},{col}##','')
    if not tables_dict:
        return re.sub(r'@@-?\d+,-?\d+##', '', text)
    
    tbs = []
    for t in tables_dict.values():
        table = []
        for row in t:
            row_c = [t[row][col].strip() for col in t[row]]
            if ''.join(row_c).strip():
                table.append(row_c)
        if not table:
            continue
        table = np.array(table)
        has_content_cols = []
        for i in range(table.shape[1]):
            if ''.join(table[:,i]).strip():
                has_content_cols.append(i)
        table = table[:, has_content_cols]
        tbs.append(md_table(table))
    text = re.sub(r'@@-?\d+,-?\d+##', '', text)
    text = re.sub('\n+', '\n', text)
    return text.rstrip('\n') + '\n' + '\n'.join(tbs)

class RAGFlowPptParser(object):
    def __init__(self):
        super().__init__()
        self.ocr = None
        self.drop_score = 0.5
        self.img_map = {}
        self.images_list = []
        self.embed_images = {}
        self.vml_pat = re.compile('vmlDrawing" Target="[^>]+drawings/([^>]+)"/>')
        self.img_pat = re.compile('image" Target="[^>]+media/([^>]+)"/>')
        self.rel_pat = re.compile('<Relationship[^>]+Type="[^>]+" Target="[^>]+"/>')
        
    def image_to_mode(self, img_blob):
        img_md5 = hashlib.md5(img_blob).hexdigest()
        if img_md5 in self.img_map:
            image_uuid = self.img_map.get(img_md5)
        else:
            self.img_map[img_md5] = uuid.uuid1()
            image_uuid = self.img_map.get(img_md5)
        if len([i for i in self.images_list if i.id == image_uuid]) == 0:
            image = ImageModel(id=image_uuid, image=img_blob, image_name=img_md5)
            self.images_list.append(image)
        return f'![](/api/image/{image_uuid})'
    
    def __extract_by_rid(self, shape):
        rids = []
        for ele in shape.element.iter():
            for v in ele.attrib.values():
                if 'rId' in v and v not in rids:
                    rids.append(v)
        texts = []
        for rid in rids:
            rel = shape.part.rels[rid]
            if 'image' in rel.reltype:
                img_blob = rel.target_part.blob
                img_name = rel.target_partname
                if img_name.endswith('.emf') or img_name.endswith('.wmf'):
                    ext = 'wmf' if img_name.endswith('.wmf') else 'emf'
                    img_blob = emf_to_png(img_blob, ext)
                text = self.image_to_mode(img_blob)
                if text not in texts:
                    texts.append(text)
            if 'diagramData' in rel.reltype:
                text = ' '.join(re.findall('<a:t>(.*?)</a:t>', rel.target_part.blob.decode()))
                texts.append(text)
            if 'vmlDrawing' in rel.reltype:
                for rel in rel.target_part.rels.values():
                    if 'image' in rel.reltype:
                        img_blob = rel.target_part.blob
                        img_name = rel.target_partname
                        if img_name.endswith('.emf') or img_name.endswith('.wmf'):
                            ext = 'wmf' if img_name.endswith('.wmf') else 'emf'
                            img_blob = emf_to_png(img_blob, ext)
                        text = self.image_to_mode(img_blob)
                        if text not in texts:
                            texts.append(text)
        return '\n'.join(texts)

    def __extract(self, shape, pos=False):
        if shape.shape_type == 13 or (hasattr(shape, 'image') and shape.image):
            try:
                img_blob = shape.image.blob
                img_name = shape.image.filename
                if img_name.endswith('.emf') or img_name.endswith('.wmf'):
                    ext = 'wmf' if img_name.endswith('.wmf') else 'emf'
                    img_blob = emf_to_png(img_blob, ext)
                img = Image.open(BytesIO(img_blob))
                img_api = self.image_to_mode(img_blob)
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                img = np.array(img)
                if self.ocr is None:
                    self.ocr = OCR()
                results = self.ocr(img)
            except Exception as e:
                print(e)
                return None
            txts = []
            for _, res in results:
                if res[1]<self.drop_score:
                    continue 
                txts.append(res[0])
            txt = ' '.join(txts)
            txt = img_api+'\n'+txt
            return txt
        
        if shape.shape_type == 19:
            tb = shape.table
            table = [[f"{tb.cell(i,j).text}" if tb.cell(i,j).text else '' 
                      for j in range(len(tb.columns))] for i in range(len(tb.rows))]
            return md_table(table)

        if shape.has_text_frame:
            if pos:
                return f"{shape.text_frame.text.strip()}@@{shape.top},{shape.left}##"
            else:
                return shape.text_frame.text

        if shape.shape_type == 6:
            texts = []
            for p in sorted(shape.shapes, key=lambda x: (x.top // 10, x.left)):
                t = self.__extract(p, pos=True)
                if t:
                    texts.append(t)
            return "\n".join(texts)
        
        return ''
        
    # def ashape_ocr(self, ashape):
    #     buffered = BytesIO()
    #     ashape.get_thumbnail().save(buffered, drawing.imaging.ImageFormat.jpeg,)
    #     if self.ocr is None:
    #         self.ocr = OCR()
    #     results = self.ocr(255-np.array(Image.open(buffered)))
    #     txts = []
    #     for _, res in results:
    #         if res[1]<self.drop_score:
    #             continue 
    #         txts.append(res[0])
    #     txt = ' '.join(txts)
    #     txt = re.sub("Evaluation\s+only.?\s+Created\s+with\s+","",txt)
    #     return txt


    def __call__(self, fnm, buffer=None, from_page=0, to_page=999999):
        if buffer is None:
            ppt = Presentation(fnm)
        else:
            ppt = Presentation(BytesIO(buffer))
                
        txts = []
        self.total_page = len(ppt.slides)
        for i, slide in enumerate(ppt.slides):
            if i < from_page:
                continue
            if i >= to_page:
                break
            texts = []
            shapes = sorted(
                    slide.shapes, key=lambda x: (x.top // 10 if x.top else 0, x.left if x.left else 0))
            for j, shape in enumerate(shapes):
                txt = self.__extract(shape)
                if shape.shape_type==6:
                    txt = build_table(txt)
                if txt.strip():
                    texts.append(txt)
                else:
                    txt = self.__extract_by_rid(shape)
                    if txt.strip():
                        texts.append(txt)
            txts.append("\n".join(texts))        
        return txts, self.images_list
